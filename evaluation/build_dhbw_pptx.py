"""
Fills the existing BA_Praesentation.pptx SAP template with TradeBridge 2.0 content.
Saves result to /Users/i589277/DHBW/BA_Praesentation_final.pptx
Run: .venv/bin/python evaluation/build_dhbw_pptx.py
"""

from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = "/Users/i589277/DHBW/BA_Präsentation.pptx"
OUT = "/Users/i589277/DHBW/BA_Präsentation_final.pptx"

# ── SAP template colours ───────────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x2A, 0x86)
SAP_BLUE   = RGBColor(0x00, 0x70, 0xF2)
GREEN      = RGBColor(0x1A, 0x7A, 0x46)
GREEN_L    = RGBColor(0xD5, 0xF0, 0xE0)
AMBER      = RGBColor(0xD4, 0x7E, 0x0F)
AMBER_L    = RGBColor(0xFF, 0xF3, 0xB8)
RED        = RGBColor(0xC0, 0x39, 0x2B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
GREY       = RGBColor(0x55, 0x55, 0x66)
LIGHT_GREY = RGBColor(0xF5, 0xF7, 0xFA)
DIVIDER_C  = RGBColor(0xD0, 0xD8, 0xE4)

prs = Presentation(SRC)

# ── Layout shortcuts ───────────────────────────────────────────
LY_DIVIDER  = prs.slide_layouts[15]   # Divider Page C
LY_TITLE    = prs.slide_layouts[17]   # Title Only
LY_TEXT     = prs.slide_layouts[18]   # Title and Text
LY_2COL     = prs.slide_layouts[19]   # Title and Text: 2 Columns
LY_3COL     = prs.slide_layouts[20]   # Title and Text: 3 Columns
LY_BLANK    = prs.slide_layouts[32]   # Blank

# ── Helper: slide insertion at position ───────────────────────

def add_slide_at(prs, layout, pos):
    """Add a new slide from layout and move it to position pos."""
    sl = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    elems = list(xml_slides)
    el = elems[-1]
    xml_slides.remove(el)
    xml_slides.insert(pos, el)
    return sl

def remove_slide(prs, index):
    """Remove slide at index."""
    xml_slides = prs.slides._sldIdLst
    rId = xml_slides[index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    del xml_slides[index]

# ── Helper: text + shapes ──────────────────────────────────────

def set_ph_text(slide, ph_idx, text, size=None, bold=None, color=None):
    """Set text of a placeholder by its idx."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            if size:  run.font.size = Pt(size)
            if bold is not None: run.font.bold = bold
            if color: run.font.color.rgb = color
            return ph
    return None

def txb(slide, text, l, t, w, h, size=12, bold=False,
        color=DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def add_para(tf, text, size=11, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align; p.space_before = space_before
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    from pptx.util import Inches as I
    sh = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

# ─────────────────────────────────────────────────────────────
# SLIDE 1 — Cover (keep as-is, already has correct title)
# ─────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────
# SLIDE 2 — Agenda (modify existing text boxes)
# ─────────────────────────────────────────────────────────────
agenda_sl = prs.slides[1]

agenda_items = [
    ("01", "Idee & Motivation",     "Warum autonome B2B-Verhandlungen? Das Problem im Tail-Spend-Einkauf."),
    ("02", "Systemarchitektur",     "Wie ist TradeBridge 2.0 aufgebaut? Agenten, Pipeline, Guardrails."),
    ("03", "KPIs & Evaluation",     "Wie messen wir Erfolg? 4 Kennzahlen, 140 Sitzungen, 14 Szenarien."),
    ("04", "Ergebnisse",            "Was zeigt die Evaluation? CSR 98,4 % — FAR 0,0 — F1 0,933."),
    ("05", "Live Demo",             "TradeBridge 2.0 in Aktion — End-to-End Verhandlung live."),
]

# Number text boxes: TextBox 27, 44, 48, 52, 4
# Content text boxes: TextBox 38, 46, 50, 54, 6
num_names   = ["TextBox 27", "TextBox 44", "TextBox 48", "TextBox 52", "TextBox 4"]
cont_names  = ["TextBox 38", "TextBox 46", "TextBox 50", "TextBox 54", "TextBox 6"]

shape_map = {sh.name: sh for sh in agenda_sl.shapes}

for i, (num, title, desc) in enumerate(agenda_items):
    # Update number box
    if num_names[i] in shape_map:
        tf = shape_map[num_names[i]].text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                r.text = num

    # Update content box: first para = title, second = desc
    if cont_names[i] in shape_map:
        tf = shape_map[cont_names[i]].text_frame
        paras = tf.paragraphs
        if len(paras) >= 1:
            for r in paras[0].runs:
                r.text = title
        if len(paras) >= 2:
            for r in paras[1].runs:
                r.text = desc
        elif len(paras) == 1:
            p2 = tf.add_paragraph()
            r = p2.add_run(); r.text = desc
            r.font.size = Pt(14)

# ─────────────────────────────────────────────────────────────
# SLIDE 3 — Divider: "01  Idee & Motivation" (modify existing)
# ─────────────────────────────────────────────────────────────
div1 = prs.slides[2]
set_ph_text(div1, 0, "01  Idee & Motivation", size=36, bold=True)

# ─────────────────────────────────────────────────────────────
# SLIDE 4 — Remove placeholder "Data points" slide
# ─────────────────────────────────────────────────────────────
remove_slide(prs, 3)
# Now slides: [Cover, Agenda, Divider01, ThankYou]  (ThankYou is now index 3)

INSERT = 3  # Insert everything before Thank You (index 3)

# ─────────────────────────────────────────────────────────────
# SLIDE 4 — Problem (Title and Text: 3 Columns)
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_3COL, INSERT); INSERT += 1
set_ph_text(sl, 0, "Das Problem: Tail-Spend-Verhandlungen im Retail", size=18, bold=True)

cols = [
    ("⏱  Manueller Aufwand",
     "Einkäufer verhandeln Standardartikel mit geringem Wert manuell — hoher Zeitaufwand, niedriger strategischer ROI.",
     "Skalierbar erst durch Automatisierung."),
    ("📉  Inkonsistente Preise",
     "Ohne formalisierte Strategie variieren Einigungspreise stark — je nach Erfahrung und Tagesform.",
     "Systematische Guardrails fehlen vollständig."),
    ("🔒  Fehlende Skalierbarkeit",
     "Steigendes Sortimentsvolumen lässt sich nicht mit wachsendem Personal lösen.",
     "KI kann 24/7 parallel verhandeln — Menschen nur sequenziell."),
]

ph_idxs = [10, 12, 13]
for ph_idx, (title, body, footer) in zip(ph_idxs, cols):
    for ph in sl.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame; tf.clear()
            # Title line
            p0 = tf.paragraphs[0]
            r0 = p0.add_run(); r0.text = title
            r0.font.size = Pt(13); r0.font.bold = True; r0.font.color.rgb = DARK_BLUE
            # Body
            p1 = tf.add_paragraph(); p1.space_before = Pt(8)
            r1 = p1.add_run(); r1.text = body
            r1.font.size = Pt(11); r1.font.color.rgb = DARK
            # Footer italic
            p2 = tf.add_paragraph(); p2.space_before = Pt(10)
            r2 = p2.add_run(); r2.text = footer
            r2.font.size = Pt(10); r2.font.italic = True; r2.font.color.rgb = GREY

# Research question box (manual text box)
rect(sl, 0.55, 6.02, 12.23, 0.88, fill=AMBER_L, line=AMBER, lw=Pt(1.5))
txb(sl, "Forschungsfrage: Laesst sich B2B-Verhandlung durch ein LLM-basiertes Multi-Agenten-System"
    " innerhalb vordefinierter Preis- und Mengenparameter automatisieren?",
    0.72, 6.08, 11.9, 0.8, size=11, italic=True, color=DARK_BLUE)

# ─────────────────────────────────────────────────────────────
# SLIDE 5 — Die Idee (Title and Text)
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_TEXT, INSERT); INSERT += 1
set_ph_text(sl, 0, "Die Idee: Zwei KI-Agenten verhandeln autonom", size=18, bold=True)

for ph in sl.placeholders:
    if ph.placeholder_format.idx == 10:
        tf = ph.text_frame; tf.clear()
        items = [
            ("🤖  Zwei LLM-Agenten",
             "Supplier-Agent und Retailer-Agent verhandeln Runde fuer Runde — jeder mit eigener Persoenlichkeit, Taktik und privaten Preislimits."),
            ("🔐  Unveraenderliche Guardrails",
             "Preisuntergrenze und Mengenkonstanten sind in der Sessionkonfiguration fixiert. Der Agent kann sie nie unterschreiten — auch nicht durch Prompt-Injection."),
            ("⚙️  Orchestrator & ZOPA",
             "Der Orchestrator berechnet die Zone of Possible Agreement, prueft jede Runde per Constraint-Validator und eskaliert bei Stagnation an den Menschen."),
            ("👤  HITL nur bei Bedarf",
             "Menschlicher Eingriff ausschliesslich bei Deadlock oder aussergewoehnlichen Grenzfaellen — alles andere laeuft vollstaendig autonom."),
        ]
        first = True
        for title, body in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(0 if first else 10)
            r = p.add_run(); r.text = title
            r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DARK_BLUE
            p2 = tf.add_paragraph(); p2.space_before = Pt(3)
            r2 = p2.add_run(); r2.text = "   " + body
            r2.font.size = Pt(11); r2.font.color.rgb = DARK
            tf.add_paragraph()

# ─────────────────────────────────────────────────────────────
# SLIDE 6 — Divider: "02  Systemarchitektur"
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_DIVIDER, INSERT); INSERT += 1
set_ph_text(sl, 0, "02  Systemarchitektur", size=36, bold=True)

# ─────────────────────────────────────────────────────────────
# SLIDE 7 — Architektur (2 Columns)
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_2COL, INSERT); INSERT += 1
set_ph_text(sl, 0, "Systemarchitektur — Komponenten & Tech Stack", size=18, bold=True)

# Column 1: Components
for ph in sl.placeholders:
    if ph.placeholder_format.idx == 10:
        tf = ph.text_frame; tf.clear()
        components = [
            ("FastAPI Backend",          "REST-API, Session-Management, ZOPA-Berechnung"),
            ("LLM-Pipeline (4 Stufen)",  "Taktik → Aspiration → Risk → Offer — sequenziell"),
            ("Constraint Validator",     "Validierung auf raw_offer; max. 3 Re-Prompts vor Clamping"),
            ("Opponent Model",           "Boulware / Conceder nach Faratin et al. (1998)"),
            ("Evaluationsmodul",         "CSR / WAA / ZU / BP — 140 Sitzungen automatisiert"),
            ("React Frontend",           "Supplier- & Retailer-Dashboard, HITL-Panel, Live-Timeline"),
        ]
        first = True
        for title, desc in components:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False; p.space_before = Pt(8)
            r = p.add_run(); r.text = title
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = DARK_BLUE
            p2 = tf.add_paragraph(); p2.space_before = Pt(2)
            r2 = p2.add_run(); r2.text = "   " + desc
            r2.font.size = Pt(10); r2.font.color.rgb = GREY

# Column 2: Tech Stack + Flow
for ph in sl.placeholders:
    if ph.placeholder_format.idx == 11:
        tf = ph.text_frame; tf.clear()
        stack = [
            ("LLM",        "Claude Sonnet 4.x  |  Temperature = 0"),
            ("Backend",    "Python 3.13  |  FastAPI  |  Pydantic v2"),
            ("Frontend",   "React 19  |  TypeScript  |  Tailwind CSS"),
            ("Protokoll",  "A2A  |  REST-basiert"),
            ("Evaluation", "14 Szenarien  |  10 Laeufe  |  140 Sitzungen"),
            ("Persistenz", "JSON-Session-Store  |  raw_offer-Feld"),
        ]
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = "Tech Stack"
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = SAP_BLUE
        for key, val in stack:
            p2 = tf.add_paragraph(); p2.space_before = Pt(7)
            r2 = p2.add_run(); r2.text = key + ":"
            r2.font.size = Pt(11); r2.font.bold = True; r2.font.color.rgb = DARK_BLUE
            p3 = tf.add_paragraph(); p3.space_before = Pt(2)
            r3 = p3.add_run(); r3.text = "   " + val
            r3.font.size = Pt(10); r3.font.color.rgb = DARK

# ─────────────────────────────────────────────────────────────
# SLIDE 8 — Divider: "03  KPIs & Evaluation"
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_DIVIDER, INSERT); INSERT += 1
set_ph_text(sl, 0, "03  KPIs & Evaluation", size=36, bold=True)

# ─────────────────────────────────────────────────────────────
# SLIDE 9 — KPIs (3 Columns)
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_3COL, INSERT); INSERT += 1
set_ph_text(sl, 0, "Vier KPIs — Was wird gemessen und warum?", size=18, bold=True)

kpi_cols = [
    [
        ("CSR — Constraint Satisfaction Rate",
         "Haelt der Agent Preis- und Mengenlimits aus eigener Kraft ein?",
         "Gemessen auf raw_offer (erstem LLM-Output, vor Korrektur).",
         "Kriterium: >= 95 %"),
        ("ZU — ZOPA Utilization",
         "Wo landet die Einigung im Verhandlungskorridor?",
         "0 = Supplier-Maximum | 0,5 = Mitte | 1 = Supplier-Minimum",
         "Diagnostisch — kein Schwellenwert"),
    ],
    [
        ("WAA — Walk-Away Accuracy",
         "Trifft der Agent die richtige Deal-/Abbruch-Entscheidung?",
         "F1 auf ZOPA/No-ZOPA-Klassifikation. FAR = 0 hartes Kriterium.",
         "Kriterium: F1 >= 0,90  |  FAR = 0,0"),
        ("BP — Business Plausibility",
         "Klingt die Verhandlung wie echte B2B-Kommunikation?",
         "Human-as-a-Judge: Argumentationsqualitaet, Konzessionslogik, Reaktionsadaequatheit.",
         "Human-as-a-Judge (kein LLM-Judge)"),
    ],
    [
        ("Hybrides Design",
         "CSR, WAA, ZU automatisiert (140 Sitzungen). BP qualitativ durch Retail-Experten.",
         "N = 10 Laeufe je Szenario — Mittelwert + Standardabweichung.",
         "14 Szenarien x 10 Laeufe = 140 Sitzungen"),
        ("Warum kein LLM-Judge fuer BP?",
         "Zheng et al. (2023): Positional Bias + Verbosity Bias.",
         "Handelsmargen, Drucktaktiken, Konzessionsmuster benoetigen echte Branchenerfahrung.",
         "Methodische Zirkularitaet vermieden"),
    ],
]

for ph_idx, col_items in zip([10, 12, 13], kpi_cols):
    for ph in sl.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame; tf.clear()
            first = True
            for title, body, method, crit in col_items:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False; p.space_before = Pt(0 if first else 12)
                r = p.add_run(); r.text = title
                r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = DARK_BLUE
                for line, col, sz in [(body, DARK, 10), (method, GREY, 9), (crit, SAP_BLUE, 10)]:
                    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
                    r2 = p2.add_run(); r2.text = line
                    r2.font.size = Pt(sz); r2.font.color.rgb = col
                    r2.font.bold = (col == SAP_BLUE)
                tf.add_paragraph()

# ─────────────────────────────────────────────────────────────
# SLIDE 10 — Divider: "04  Ergebnisse"
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_DIVIDER, INSERT); INSERT += 1
set_ph_text(sl, 0, "04  Ergebnisse", size=36, bold=True)

# ─────────────────────────────────────────────────────────────
# SLIDE 11 — Results Overview (Title Only + manual shapes)
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_TITLE, INSERT); INSERT += 1
set_ph_text(sl, 0, "Ergebnisse — 140 Sitzungen, 14 Szenarien, 10 Laeufe", size=18, bold=True)

# 4 big metric cards
metrics = [
    ("CSR",      "98,4 %",   "± 1,6 %",   ">= 95 %  erfuellt",   GREEN,     GREEN_L),
    ("WAA  F1",  "0,933",    "",          ">= 0,90  erfuellt",   GREEN,     GREEN_L),
    ("WAA  FAR", "0,000",    "",          "= 0,0    erfuellt",   GREEN,     GREEN_L),
    ("BP",       "8 / 14",   "ueberwiegend / hoch", "Human-as-a-Judge", SAP_BLUE, RGBColor(0xD6,0xEA,0xF8)),
]
for i, (label, val, sub, crit, col, bg) in enumerate(metrics):
    l = 0.55 + i * 3.2
    rect(sl, l, 1.1, 3.05, 2.4, fill=DARK_BLUE)
    rect(sl, l, 1.1, 3.05, 0.06, fill=col)
    txb(sl, label, l+0.12, 1.2,  2.8, 0.35, size=11,
        color=RGBColor(0x90,0xB8,0xD8))
    txb(sl, val,   l+0.1,  1.55, 2.85, 0.7,  size=28, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    if sub:
        txb(sl, sub, l+0.1, 2.2, 2.85, 0.28, size=9,
            color=RGBColor(0x80,0xA0,0xC0), align=PP_ALIGN.CENTER)
    txb(sl, crit, l+0.1, 2.5, 2.85, 0.35, size=10, bold=True,
        color=col, align=PP_ALIGN.CENTER)

# Three insight boxes
insights = [
    (GREEN,    "Null BATNA-Verletzungen",
     "FAR = 0,0 — in keiner der 140 Sitzungen ein Deal unterhalb des Reservationspreises."),
    (AMBER,    "10 FN: verlorene Deals",
     "Narrow-ZOPA-Szenarien (S04-S06): enges Preisfenster fuehrt zu 20-40 % Fehlquote."),
    (SAP_BLUE, "Volume-Leverage wirkt",
     "Mengenhebel in 10/10 Laeufen korrekt eingesetzt (S13). ZU = 0,618 — Haendler-Vorteil."),
]
for i, (col, title, desc) in enumerate(insights):
    l = 0.55 + i * 4.28
    rect(sl, l, 3.72, 4.06, 1.95, fill=WHITE, line=col, lw=Pt(2))
    txb(sl, title, l+0.15, 3.82, 3.75, 0.38, size=12, bold=True, color=DARK_BLUE)
    txb(sl, desc,  l+0.15, 4.25, 3.75, 1.25, size=10, color=GREY, italic=True)

# Confusion matrix mini
rect(sl, 0.55, 5.85, 6.2, 1.35, fill=RGBColor(0xF5,0xF7,0xFA), line=DIVIDER_C, lw=Pt(1))
txb(sl, "Konfusionsmatrix (N = 120)", 0.7, 5.92, 5.8, 0.3, size=10, bold=True, color=SAP_BLUE)
for ci, (lbl, col) in enumerate([("TP=70",GREEN),("FP=0",GREEN),("FN=10",AMBER),("TN=40",GREEN)]):
    cl = 0.7 + (ci % 2) * 1.55
    ct = 6.3 + (ci // 2) * 0.52
    rect(sl, cl, ct, 1.4, 0.45, fill=col)
    txb(sl, lbl, cl+0.05, ct+0.06, 1.3, 0.32, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(sl, "Precision=1,000  |  Recall=0,875  |  F1=0,933  |  FAR=0,000  |  FWR=0,125",
    7.0, 6.55, 6.0, 0.38, size=11, bold=True, color=DARK_BLUE)
txb(sl, "FWR=0,125 (verlorene Deals) liegt unter Toleranzgrenze 0,25.",
    7.0, 6.98, 6.0, 0.3, size=10, color=GREY, italic=True)

# ─────────────────────────────────────────────────────────────
# SLIDE 12 — ZU + BP (2 Columns)
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_2COL, INSERT); INSERT += 1
set_ph_text(sl, 0, "ZOPA Utilization & Business Plausibility", size=18, bold=True)

# ZU column
for ph in sl.placeholders:
    if ph.placeholder_format.idx == 10:
        tf = ph.text_frame; tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = "ZOPA Utilization — Oe 0,303 +/- 0,043"
        r0.font.size = Pt(13); r0.font.bold = True; r0.font.color.rgb = SAP_BLUE

        zu_rows = [
            ("S01-S03", "Wide ZOPA",       "0,219 / 0,160 / 0,312", "Supplier-dominant"),
            ("S04-S06", "Narrow ZOPA",     "0,094 / 0,286 / 0,183", "Supplier-dominant"),
            ("S11-S12", "Asymmetrisch",    "0,293 / 0,279",          "Supplier-dominant"),
            ("S13",     "Volume Leverage", "0,618",                   "Haendler-dominant — Hebel wirkt"),
            ("S14",     "Volume Leverage", "0,517",                   "Ausgeglichen"),
        ]
        for sid, cat, vals, interp in zu_rows:
            p = tf.add_paragraph(); p.space_before = Pt(8)
            r = p.add_run(); r.text = f"{sid}  ({cat})"
            r.font.size = Pt(11); r.font.bold = True
            r.font.color.rgb = SAP_BLUE if "Hebel" in interp else DARK_BLUE
            p2 = tf.add_paragraph(); p2.space_before = Pt(2)
            r2 = p2.add_run(); r2.text = f"   ZU: {vals}  |  {interp}"
            r2.font.size = Pt(10); r2.font.color.rgb = GREY

        tf.add_paragraph()
        p_note = tf.add_paragraph(); p_note.space_before = Pt(10)
        r_note = p_note.add_run()
        r_note.text = "Supplier-Bias systembedingt (Ankerstrategie). Korrigierbar ueber initial_anchor-Parameter."
        r_note.font.size = Pt(10); r_note.font.italic = True; r_note.font.color.rgb = GREY

# BP column
for ph in sl.placeholders:
    if ph.placeholder_format.idx == 11:
        tf = ph.text_frame; tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = "Business Plausibility — Einstufungen"
        r0.font.size = Pt(13); r0.font.bold = True; r0.font.color.rgb = SAP_BLUE

        bp_items = [
            ("Hoch plausibel (1)",          "S13 — Volume-Leverage-Mechanismus",   GREEN),
            ("Ueberwiegend plausibel (7)",   "S01-S04, S12, S14",                   DARK_BLUE),
            ("Eingeschraenkt plausibel (6)", "S05-S11",                             AMBER),
            ("Nicht plausibel (0)",          "—",                                   GREY),
        ]
        for label, examples, col in bp_items:
            p = tf.add_paragraph(); p.space_before = Pt(10)
            r = p.add_run(); r.text = label
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = col
            p2 = tf.add_paragraph(); p2.space_before = Pt(3)
            r2 = p2.add_run(); r2.text = "   " + examples
            r2.font.size = Pt(10); r2.font.color.rgb = GREY

        tf.add_paragraph()
        p_note = tf.add_paragraph(); p_note.space_before = Pt(10)
        r_note = p_note.add_run()
        r_note.text = "Hauptschwaeche: keine aktive Walk-Away-Kommunikation in No-ZOPA-Szenarien (S07-S10)."
        r_note.font.size = Pt(10); r_note.font.italic = True; r_note.font.color.rgb = GREY

# ─────────────────────────────────────────────────────────────
# SLIDE 13 — Divider: "05  Live Demo"
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_DIVIDER, INSERT); INSERT += 1
set_ph_text(sl, 0, "05  Live Demo", size=36, bold=True)

# ─────────────────────────────────────────────────────────────
# SLIDE 14 — Demo Cue (Title and Text)
# ─────────────────────────────────────────────────────────────
sl = add_slide_at(prs, LY_TEXT, INSERT); INSERT += 1
set_ph_text(sl, 0, "TradeBridge 2.0 — End-to-End Demo", size=18, bold=True)

for ph in sl.placeholders:
    if ph.placeholder_format.idx == 10:
        tf = ph.text_frame; tf.clear()
        steps = [
            ("1.  Supplier-Dashboard",
             "Produkt auswaehlen (z.B. Bosch GSR 18V). Angebotspreis + private Preisuntergrenze setzen. Angebot an Retailer senden."),
            ("2.  Retailer-Dashboard",
             "Angebot im Posteingang empfangen. Eigene Max-Grenze setzen (vertraulich). Auto-Verhandlung starten."),
            ("3.  KI-Agenten verhandeln",
             "Runde fuer Runde: Taktik, Begruendung, Konzession. Live-Timeline zeigt Preisverlauf + AgentReasoning."),
            ("4.  Abschluss / HITL",
             "Einigung → beide Parteien bestaetigen. Stagnation → HITL-Eskalation. Kein Deal → Abbruch mit Status-Report."),
        ]
        first = True
        for title, desc in steps:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False; p.space_before = Pt(0 if first else 12)
            r = p.add_run(); r.text = title
            r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DARK_BLUE
            p2 = tf.add_paragraph(); p2.space_before = Pt(4)
            r2 = p2.add_run(); r2.text = "   " + desc
            r2.font.size = Pt(11); r2.font.color.rgb = DARK
            tf.add_paragraph()

# ─────────────────────────────────────────────────────────────
# SLIDE 15 — Thank You (keep, now last)
# ─────────────────────────────────────────────────────────────
# No changes needed — already in correct position

# ── Save ──────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved -> {OUT}")
print(f"Total slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides):
    texts = [sh.text_frame.text.strip()[:50]
             for sh in sl.shapes if sh.has_text_frame
             and sh.text_frame.text.strip()]
    print(f"  Slide {i+1}: {texts[0] if texts else '(no text)'}")
