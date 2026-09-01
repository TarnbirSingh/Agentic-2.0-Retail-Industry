"""
Creates TradeBridge 2.0 evaluation slide deck (presentation_tradebridge.pptx).
Run: .venv/bin/python evaluation/create_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x23, 0x3C)   # backgrounds / titles
MID_BLUE    = RGBColor(0x2D, 0x4A, 0x8A)   # accent
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF7)   # soft background boxes
GREEN       = RGBColor(0x2E, 0x86, 0x48)   # success / fulfilled
AMBER       = RGBColor(0xD4, 0x7E, 0x0F)   # partial
RED         = RGBColor(0xC0, 0x39, 0x2B)   # warning
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF4, 0xF6, 0xF9)
DARK_GREY   = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def para(tf, text, size=14, bold=False, color=DARK_GREY,
         align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.35, fill=DARK_BLUE)
    txbox(slide, title,   0.4, 0.12, 10, 0.7,  size=30, bold=True,  color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.82, 10, 0.45, size=15, color=RGBColor(0xA8,0xC4,0xE8))

def footer(slide, text="TradeBridge 2.0  |  Bachelorarbeit  |  DHBW Mannheim / SAP SE  |  2026"):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=DARK_BLUE)
    txbox(slide, text, 0.3, 7.17, 12, 0.28, size=9, color=RGBColor(0x90,0xA8,0xC8))

def kpi_box(slide, l, t, w, h, label, value, sub, bg=LIGHT_BLUE,
            ok=True, value_color=None):
    rect(slide, l, t, w, h, fill=bg, line=MID_BLUE, line_w=Pt(1.2))
    txbox(slide, label, l+0.12, t+0.08, w-0.2, 0.32, size=11, bold=True,
          color=MID_BLUE)
    vc = value_color or (GREEN if ok else RED)
    txbox(slide, value, l+0.1, t+0.38, w-0.15, 0.55, size=26, bold=True, color=vc)
    txbox(slide, sub, l+0.12, t+0.95, w-0.2, 0.38, size=10, color=DARK_GREY,
          italic=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
rect(sl, 0, 2.8, 13.33, 0.06, fill=MID_BLUE)

txbox(sl, "TradeBridge 2.0", 1.0, 1.2, 11, 1.1,
      size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "LLM-basiertes Multi-Agenten-System für B2B-Verhandlungen",
      1.0, 2.35, 11, 0.55, size=20, color=RGBColor(0xA8,0xC4,0xE8),
      align=PP_ALIGN.CENTER)
txbox(sl, "Evaluation — Bachelorarbeit",
      1.0, 3.05, 11, 0.5, size=18, color=RGBColor(0x90,0xA8,0xC8),
      align=PP_ALIGN.CENTER)
txbox(sl, "Tarnbir Singh  ·  DHBW Mannheim / SAP SE  ·  2026",
      1.0, 3.7, 11, 0.45, size=14, color=RGBColor(0x70,0x90,0xB0),
      align=PP_ALIGN.CENTER)
txbox(sl, "140 Verhandlungssitzungen  ·  14 Szenarien  ·  10 Läufe",
      1.0, 4.35, 11, 0.45, size=13, color=RGBColor(0x60,0x80,0xA8),
      align=PP_ALIGN.CENTER, italic=True)
footer(sl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — System Architecture
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(sl, "Systemarchitektur", "TradeBridge 2.0 — Überblick")

# Three agent boxes
agent_tops = [1.6, 1.6, 1.6]
agent_lefts = [0.4, 5.0, 9.5]
agent_labels = ["Supplier\nAgent", "Orchestrator", "Retailer\nAgent"]
agent_colors = [MID_BLUE, DARK_BLUE, RGBColor(0x1E, 0x6B, 0x4A)]
agent_icons  = ["🤖", "⚙️", "🤖"]

for i in range(3):
    rect(sl, agent_lefts[i], agent_tops[i], 3.2, 1.9,
         fill=agent_colors[i], line=WHITE, line_w=Pt(1.5))
    txbox(sl, agent_icons[i], agent_lefts[i]+1.3, agent_tops[i]+0.1,
          0.7, 0.5, size=22, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, agent_labels[i], agent_lefts[i]+0.1, agent_tops[i]+0.6,
          3.0, 0.7, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows between boxes
txbox(sl, "←  Angebot / Gegenangebot  →",
      3.7, 2.25, 1.2, 0.6, size=9, color=DARK_GREY, align=PP_ALIGN.CENTER)
txbox(sl, "←  Angebot / Gegenangebot  →",
      8.25, 2.25, 1.2, 0.6, size=9, color=DARK_GREY, align=PP_ALIGN.CENTER)

# Feature boxes below agents
features = [
    ["LLM-Pipeline (4 Schritte)", "Constraint Validator", "Retry-Mechanismus (max 3×)", "raw_offer Persistenz"],
    ["Session Management", "ZOPA-Metadaten", "HITL-Eskalation", "Evaluationsmodul"],
    ["LLM-Pipeline (4 Schritte)", "Opponent Model", "BATNA-Guardrails", "Walk-Away-Logik"],
]
for i, feats in enumerate(features):
    rect(sl, agent_lefts[i], 3.65, 3.2, 2.15, fill=WHITE,
         line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
    tb = sl.shapes.add_textbox(
        Inches(agent_lefts[i]+0.12), Inches(3.75),
        Inches(2.96), Inches(1.9))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    for j, f in enumerate(feats):
        pp = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
        pp.text = f"  ·  {f}"
        pp.font.size = Pt(11)
        pp.font.color.rgb = DARK_GREY

# Flow label
txbox(sl, "Claude Sonnet  ·  Temperature = 0  ·  A2A-Protokoll",
      3.0, 6.0, 7.3, 0.45, size=12, color=MID_BLUE,
      align=PP_ALIGN.CENTER, italic=True)
footer(sl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Evaluation Design
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(sl, "Evaluationsdesign", "Hybrid: quantitativ + qualitativ — 140 Sitzungen")

# Left column — Design info
rect(sl, 0.35, 1.5, 5.8, 5.3, fill=WHITE, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
txbox(sl, "Aufbau", 0.55, 1.6, 5.4, 0.4, size=13, bold=True, color=MID_BLUE)

tb = sl.shapes.add_textbox(Inches(0.55), Inches(2.1), Inches(5.4), Inches(4.4))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
items = [
    ("14 Szenarien", "Wide / Narrow / No-ZOPA, Near-Miss, Asymmetrisch, Volume-Leverage"),
    ("10 Läufe je Szenario", "N = 140 Verhandlungssitzungen, ca. 1.300 Runden"),
    ("Fixierte Parameter", "Temperature = 0, deterministische Agentenpersönlichkeiten"),
    ("Multi-Run-Statistik", "Mittelwert ± Standardabweichung pro KPI"),
]
first = True
for title, desc in items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run()
    run.text = f"▸  {title}"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = f"    {desc}"
    run2.font.size = Pt(11)
    run2.font.italic = True
    run2.font.color.rgb = DARK_GREY
    tf.add_paragraph()

# Right column — KPI overview
rect(sl, 6.5, 1.5, 6.5, 5.3, fill=WHITE, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
txbox(sl, "Vier KPIs", 6.7, 1.6, 6.0, 0.4, size=13, bold=True, color=MID_BLUE)

kpi_data = [
    ("CSR", "Constraint Satisfaction Rate", "Hält der Agent seine Preislimits ein?", GREEN),
    ("WAA", "Walk-Away Accuracy", "Entscheidet er korrekt: Deal oder Abbruch?", MID_BLUE),
    ("ZU",  "ZOPA Utilization", "Wo landet die Einigung im Korridor?", AMBER),
    ("BP",  "Business Plausibility", "Klingt es wie echte B2B-Verhandlung?", RGBColor(0x6A, 0x1B, 0x9A)),
]

for i, (short, name, desc, col) in enumerate(kpi_data):
    top = 2.15 + i * 1.15
    rect(sl, 6.6, top, 6.2, 0.95, fill=LIGHT_BLUE,
         line=col, line_w=Pt(2))
    txbox(sl, short, 6.75, top+0.06, 0.7, 0.4, size=18, bold=True, color=col)
    txbox(sl, name, 7.55, top+0.08, 5.1, 0.35, size=12, bold=True, color=DARK_BLUE)
    txbox(sl, desc, 7.55, top+0.48, 5.1, 0.35, size=10, color=DARK_GREY, italic=True)

footer(sl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — CSR Results
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(sl, "CSR — Constraint Satisfaction Rate",
           "Hält das System seine Preis- und Mengengrenzen aus eigener Kraft ein?")

# Big KPI card
rect(sl, 0.35, 1.55, 3.8, 2.0, fill=DARK_BLUE)
txbox(sl, "98,4 %", 0.45, 1.65, 3.6, 1.0, size=44, bold=True, color=GREEN,
      align=PP_ALIGN.CENTER)
txbox(sl, "± 1,6 %  ·  Ø über 10 Läufe", 0.45, 2.65, 3.6, 0.45,
      size=13, color=RGBColor(0xA8,0xC4,0xE8), align=PP_ALIGN.CENTER)
txbox(sl, "✓  Kriterium ≥ 95 % erfüllt", 0.45, 3.05, 3.6, 0.38,
      size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# What it means
tb = sl.shapes.add_textbox(Inches(4.4), Inches(1.6), Inches(8.6), Inches(1.8))
tb.word_wrap = True
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Was bedeutet das?"
run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = MID_BLUE
items2 = [
    "In 19 von 20 Runden hält das LLM die Constraints ohne Korrektursystem ein.",
    "Gemessen ausschließlich auf raw_offer — dem unkorrigierten ersten LLM-Output.",
    "Keine einzige Preis-Constraint-Verletzung in 140 Sitzungen."
]
for it in items2:
    pp = tf.add_paragraph()
    run2 = pp.add_run(); run2.text = f"  ·  {it}"
    run2.font.size = Pt(11); run2.font.color.rgb = DARK_GREY

# Per-scenario bar chart (horizontal bars via shapes)
rect(sl, 0.35, 3.75, 12.6, 3.1, fill=WHITE,
     line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
txbox(sl, "CSR pro Szenario (Mittelwert, alle 10 Läufe)",
      0.5, 3.82, 8, 0.32, size=11, bold=True, color=MID_BLUE)

scenarios = ["S01","S02","S03","S04","S05","S06",
             "S07","S08","S09","S10","S11","S12","S13","S14"]
csr_vals  = [0.975,0.950,0.917,1.0,1.0,1.0,
             1.0,1.0,1.0,1.0,0.946,0.988,1.0,1.0]

bar_max_w = 9.6
bar_h     = 0.165
bar_top0  = 4.25
bar_left  = 2.25

for i, (sc_id, val) in enumerate(zip(scenarios, csr_vals)):
    top = bar_top0 + i * 0.195
    color = GREEN if val >= 0.95 else AMBER
    txbox(sl, sc_id, 0.45, top, 0.7, bar_h+0.02, size=9, color=DARK_GREY)
    w = bar_max_w * val
    rect(sl, bar_left, top, w, bar_h, fill=color)
    txbox(sl, f"{val:.3f}", bar_left + w + 0.05, top, 0.7, bar_h,
          size=8, color=DARK_GREY)

# Threshold line (visual approximation)
rect(sl, bar_left + bar_max_w * 0.95, 4.2, 0.02, 2.75,
     fill=RED, line=RED, line_w=Pt(1))
txbox(sl, "95 %", bar_left + bar_max_w * 0.95 - 0.15, 4.05,
      0.5, 0.2, size=8, color=RED)

footer(sl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — WAA Results
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(sl, "WAA — Walk-Away Accuracy",
           "Trifft das System die richtige Deal/Abbruch-Entscheidung?")

# Confusion matrix
rect(sl, 0.35, 1.55, 5.9, 3.1, fill=WHITE,
     line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
txbox(sl, "Konfusionsmatrix (120 klassifizierte Sitzungen, N=10 Läufe)",
      0.5, 1.62, 5.6, 0.38, size=11, bold=True, color=MID_BLUE)

# Matrix headers
txbox(sl, "ZOPA vorhanden", 2.1, 2.1, 2.1, 0.3, size=10, bold=True,
      color=DARK_BLUE, align=PP_ALIGN.CENTER)
txbox(sl, "Keine ZOPA", 4.2, 2.1, 1.8, 0.3, size=10, bold=True,
      color=DARK_BLUE, align=PP_ALIGN.CENTER)
txbox(sl, "Deal", 0.5, 2.6, 1.5, 0.5, size=12, bold=True,
      color=DARK_BLUE, align=PP_ALIGN.RIGHT)
txbox(sl, "Abbruch", 0.5, 3.3, 1.5, 0.5, size=12, bold=True,
      color=DARK_BLUE, align=PP_ALIGN.RIGHT)

cell_data = [
    (2.1, 2.5, "TP = 70", GREEN),
    (4.2, 2.5, "FP = 0",  GREEN),
    (2.1, 3.2, "FN = 10", AMBER),
    (4.2, 3.2, "TN = 40", GREEN),
]
for (cl, ct, txt, col) in cell_data:
    rect(sl, cl, ct, 1.9, 0.65, fill=col)
    txbox(sl, txt, cl+0.05, ct+0.1, 1.8, 0.45, size=16, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)

# KPI cards on the right
kpis_waa = [
    ("F₁ = 0,933",  "≥ 0,90  ✓", GREEN),
    ("FAR = 0,000", "= 0,0   ✓", GREEN),
    ("FWR = 0,125", "≤ 0,25  ✓", GREEN),
    ("Precision = 1,000", "kein BATNA-Fehler", GREEN),
]
for i, (val, crit, col) in enumerate(kpis_waa):
    l = 6.5 + (i % 2) * 3.35
    t = 1.58 + (i // 2) * 1.6
    rect(sl, l, t, 3.1, 1.35, fill=DARK_BLUE)
    txbox(sl, val, l+0.1, t+0.1, 2.9, 0.65, size=21, bold=True,
          color=col, align=PP_ALIGN.CENTER)
    txbox(sl, crit, l+0.1, t+0.78, 2.9, 0.45, size=12,
          color=RGBColor(0xA8,0xC4,0xE8), align=PP_ALIGN.CENTER)

# Insight boxes
rect(sl, 0.35, 4.75, 6.0, 2.1, fill=WHITE,
     line=GREEN, line_w=Pt(2))
txbox(sl, "✓  Kein einziger BATNA-verletzender Deal",
      0.5, 4.85, 5.7, 0.4, size=12, bold=True, color=GREEN)
txbox(sl, "FAR = 0,0 über alle 140 Sitzungen — in keinem der\n"
          "No-ZOPA-Szenarien (S07–S10) wurde ein Deal abgeschlossen.",
      0.5, 5.3, 5.7, 0.75, size=11, color=DARK_GREY)
txbox(sl, "Asymmetrische Fehlerlogik:\nFP (BATNA-Verletzung) = absolut inakzeptabel\n"
          "FN (verlorener Deal)  = tolerierbar für PoC",
      0.5, 6.08, 5.7, 0.65, size=10, color=DARK_GREY, italic=True)

rect(sl, 6.5, 4.75, 6.5, 2.1, fill=WHITE,
     line=AMBER, line_w=Pt(2))
txbox(sl, "⚠  10 FN-Fälle — verlorene Deals bei vorhandener ZOPA",
      6.65, 4.85, 6.2, 0.4, size=12, bold=True, color=AMBER)
txbox(sl, "S04 (2×), S05 (3×), S06 (4×), S13 Lauf 8 (1×)\n"
          "9× max_rounds_reached  ·  1× autonomer Walk-Away (S06 Lauf 8)\n"
          "Ursache: enges ZOPA-Fenster + inkonsistentes Ankerverhalten",
      6.65, 5.3, 6.2, 0.9, size=11, color=DARK_GREY)
txbox(sl, "Run-Level: Ø F₁ = 0,896 ± 0,082 — einzelne Läufe unter Schwellenwert",
      6.65, 6.2, 6.2, 0.45, size=10, color=DARK_GREY, italic=True)

footer(sl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — ZU Results
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(sl, "ZU — ZOPA Utilization",
           "Wo landet die Einigung im Verhandlungskorridor?")

# Formula box
rect(sl, 0.35, 1.55, 5.5, 1.2, fill=DARK_BLUE)
txbox(sl, "ZU = (Retailer-Max − Einigungspreis) ÷ (Retailer-Max − Supplier-Min)",
      0.5, 1.68, 5.2, 0.45, size=13, bold=True, color=WHITE)
txbox(sl, "0,0 = Supplier gewinnt alles  ·  0,5 = Mitte  ·  1,0 = Retailer gewinnt alles",
      0.5, 2.15, 5.2, 0.4, size=11, color=RGBColor(0xA8,0xC4,0xE8))

# Overall result
rect(sl, 6.1, 1.55, 6.85, 1.2, fill=DARK_BLUE)
txbox(sl, "Gesamt-ZU über 83 Einigungen:", 6.25, 1.6, 6.5, 0.35,
      size=12, color=RGBColor(0xA8,0xC4,0xE8))
txbox(sl, "0,303 ± 0,043", 6.25, 1.9, 4.0, 0.65, size=28, bold=True,
      color=AMBER, align=PP_ALIGN.LEFT)
txbox(sl, "→ Supplier-dominant", 9.8, 1.95, 2.9, 0.5, size=15,
      color=AMBER, bold=True)

# ZU spectrum bar
rect(sl, 0.35, 2.95, 12.6, 0.5, fill=WHITE,
     line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
txbox(sl, "0,0", 0.4, 3.0, 0.5, 0.35, size=10, color=DARK_GREY)
txbox(sl, "0,40", 5.2, 3.0, 0.6, 0.35, size=10, color=DARK_GREY)
txbox(sl, "0,60", 7.8, 3.0, 0.6, 0.35, size=10, color=DARK_GREY)
txbox(sl, "1,0", 12.6, 3.0, 0.5, 0.35, size=10, color=DARK_GREY)
# supplier band
rect(sl, 0.7, 3.05, 5.6, 0.35, fill=RGBColor(0xFF,0xE0,0xB2))
txbox(sl, "Supplier-dominant  < 0,40", 0.9, 3.07, 5.0, 0.3, size=9, color=DARK_GREY)
# balanced
rect(sl, 6.3, 3.05, 2.8, 0.35, fill=RGBColor(0xC8,0xE6,0xC9))
txbox(sl, "Ausgeglichen", 6.5, 3.07, 2.4, 0.3, size=9, color=DARK_GREY)
# retailer
rect(sl, 9.1, 3.05, 3.85, 0.35, fill=LIGHT_BLUE)
txbox(sl, "Retailer-dominant  > 0,60", 9.2, 3.07, 3.5, 0.3, size=9, color=DARK_GREY)
# marker for 0.303
rect(sl, 0.7 + (12.6-0.7)*0.303, 2.88, 0.04, 0.65,
     fill=RED, line=RED, line_w=Pt(1.5))
txbox(sl, "Ø 0,303", 0.7 + (12.6-0.7)*0.303 + 0.06, 2.9,
      0.8, 0.25, size=9, bold=True, color=RED)

# Per-scenario table
rect(sl, 0.35, 3.6, 12.6, 3.2, fill=WHITE,
     line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
txbox(sl, "ZU pro Szenario (nur Sitzungen mit Einigung, N = 83)",
      0.5, 3.67, 8, 0.32, size=11, bold=True, color=MID_BLUE)

zu_data = [
    ("S01 Wide ZOPA",      0.219, 0.150, "Supplier-dominant"),
    ("S02 Wide ZOPA",      0.160, 0.107, "Supplier-dominant"),
    ("S03 Wide ZOPA",      0.312, 0.241, "Supplier-dominant"),
    ("S04 Narrow ZOPA",    0.094, 0.157, "Supplier-dominant"),
    ("S05 Narrow ZOPA",    0.286, 0.488, "Supplier-dominant"),
    ("S06 Narrow ZOPA",    0.183, 0.402, "Supplier-dominant"),
    ("S11 Asymmetrisch",   0.293, 0.162, "Supplier-dominant"),
    ("S12 Asymmetrisch",   0.279, 0.303, "Supplier-dominant"),
    ("S13 Volume Leverage",0.618, 0.176, "Händler-dominant ← Hebel wirkt"),
    ("S14 Volume Leverage",0.517, 0.189, "Ausgeglichen"),
]
col1_w = 2.5
bar_max = 5.8
bar_start = 3.0
for i, (name, mean, std, cat) in enumerate(zu_data):
    top = 4.1 + i * 0.26
    txbox(sl, name, 0.5, top, col1_w, 0.24, size=9, color=DARK_GREY)
    bar_w = bar_max * mean
    col = GREEN if mean >= 0.5 else (RGBColor(0xFF,0xA0,0x00) if mean >= 0.4 else RGBColor(0xD4,0x7E,0x0F))
    if "Händler" in cat:
        col = MID_BLUE
    rect(sl, bar_start, top+0.02, bar_w, 0.2, fill=col)
    txbox(sl, f"{mean:.3f} ± {std:.3f}", bar_start + bar_w + 0.05,
          top, 1.5, 0.22, size=8, color=DARK_GREY)
    if "Hebel" in cat:
        txbox(sl, "✓ Hebel wirkt", bar_start + bar_max + 1.65,
              top, 1.5, 0.22, size=9, bold=True, color=GREEN)

footer(sl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — BP Results
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
header_bar(sl, "BP — Business Plausibility",
           "Human-as-a-Judge: 14 Szenarien, 140 Sitzungen, ~1.300 Runden")

# BP rating distribution
bp_cats = [
    ("Hoch plausibel",          1,  GREEN,                   "S13"),
    ("Überwiegend plausibel",   7,  MID_BLUE,                "S01 S02 S03 S04 S12 S14"),
    ("Eingeschränkt plausibel", 6,  AMBER,                   "S05 S06 S07 S08 S09 S10 S11"),
    ("Nicht plausibel",         0,  RGBColor(0xEE,0xEE,0xEE),"—"),
]
for i, (label, n, col, examples) in enumerate(bp_cats):
    t = 1.6 + i * 1.28
    rect(sl, 0.35, t, 7.0, 1.1, fill=WHITE,
         line=col, line_w=Pt(2.5))
    rect(sl, 0.35, t, 1.1, 1.1, fill=col)
    txbox(sl, str(n), 0.45, t+0.18, 0.9, 0.65, size=32, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, label, 1.55, t+0.08, 4.5, 0.4, size=14, bold=True, color=DARK_BLUE)
    txbox(sl, examples, 1.55, t+0.55, 5.6, 0.4, size=10, color=DARK_GREY, italic=True)

# Right — key findings
rect(sl, 7.6, 1.55, 5.38, 5.2, fill=WHITE,
     line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
txbox(sl, "Zentrale Befunde", 7.75, 1.62, 5.0, 0.38, size=13, bold=True, color=MID_BLUE)

findings = [
    ("✓", "Volume-Leverage funktioniert",
     "Volumenhebel in allen 10 Läufen korrekt eingesetzt (S13). Konsistentestes Stärke-Muster im gesamten Test.", GREEN),
    ("✓", "Keine aggressiven Formulierungen",
     "Deeskaltiver Grundton in 140/140 Sitzungen. Geeignet für automatisierten Geschäftsverkehr.", GREEN),
    ("⚠", "Fehlende Walk-Away-Sprache",
     "No-ZOPA-Szenarien: 15 Runden lang 'win-win'-Formulierungen ohne aktiven Abbruch.", AMBER),
    ("⚠", "Generische Argumentation",
     "Kein Szenario referenziert Saisonalität, Wettbewerbspreise oder sortimentsspezifische Margen.", AMBER),
    ("⚠", "S11: kein Dimension-Shift",
     "MOQ-Konflikt (300 vs. 500 Einheiten) wird in keinem der 10 Läufe explizit benannt.", RED),
]
for i, (icon, title, desc, col) in enumerate(findings):
    top = 2.1 + i * 0.94
    txbox(sl, icon, 7.7, top, 0.35, 0.4, size=15, bold=True, color=col)
    txbox(sl, title, 8.1, top, 4.7, 0.38, size=12, bold=True, color=DARK_BLUE)
    txbox(sl, desc, 8.1, top+0.4, 4.7, 0.42, size=10, color=DARK_GREY, italic=True)

footer(sl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Summary / Gesamturteil
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)

txbox(sl, "Gesamturteil", 0.5, 0.2, 12, 0.7, size=34, bold=True, color=WHITE)
rect(sl, 0.35, 1.0, 12.6, 0.05, fill=MID_BLUE)

summary_kpis = [
    ("CSR",  "98,4 %",   "± 1,6 %",    "Kriterium ≥ 95 %",      "erfüllt ✓", GREEN),
    ("WAA\nF₁", "0,933", "",           "Kriterium ≥ 0,90",      "erfüllt ✓", GREEN),
    ("WAA\nFAR", "0,000","",           "hartes Kriterium = 0",  "erfüllt ✓", GREEN),
    ("ZU",   "0,303",    "± 0,043",    "diagnostisch (kein SW)","Supplier-dominant", AMBER),
    ("BP",   "8 / 14",   "",           "überwiegend / hoch",    "positiv ✓", MID_BLUE),
]

for i, (kpi, val, pm, crit, verdict, col) in enumerate(summary_kpis):
    l = 0.35 + i * 2.52
    rect(sl, l, 1.15, 2.42, 2.7, fill=RGBColor(0x22,0x2D,0x50))
    txbox(sl, kpi, l+0.1, 1.25, 2.2, 0.55, size=14, bold=True,
          color=RGBColor(0xA8,0xC4,0xE8), align=PP_ALIGN.CENTER)
    txbox(sl, val, l+0.1, 1.78, 2.2, 0.7, size=28, bold=True,
          color=col, align=PP_ALIGN.CENTER)
    if pm:
        txbox(sl, pm, l+0.1, 2.45, 2.2, 0.3, size=11,
              color=RGBColor(0x80,0xA0,0xC0), align=PP_ALIGN.CENTER)
    txbox(sl, crit, l+0.1, 2.78, 2.2, 0.35, size=10, italic=True,
          color=RGBColor(0x80,0xA0,0xC0), align=PP_ALIGN.CENTER)
    txbox(sl, verdict, l+0.1, 3.12, 2.2, 0.5, size=12, bold=True,
          color=col, align=PP_ALIGN.CENTER)

# Conclusion text
rect(sl, 0.35, 4.05, 12.6, 1.55, fill=RGBColor(0x22,0x2D,0x50))
txbox(sl, "Fazit",
      0.55, 4.12, 2.0, 0.38, size=13, bold=True,
      color=RGBColor(0xA8,0xC4,0xE8))
txbox(sl,
      "Alle drei quantitativen Erfolgskriterien erfüllt. Das System schließt keinen einzigen BATNA-verletzenden Deal "
      "und erkennt No-ZOPA-Situationen zuverlässig. In 8 von 14 Szenarien ist das Verhandlungsverhalten "
      "kommunikativ plausibel.",
      0.55, 4.52, 12.2, 0.65, size=13, color=WHITE)
txbox(sl,
      "Die Forschungsfrage ist positiv zu beantworten: Ein LLM-basiertes Multi-Agenten-System kann standardisierte "
      "B2B-Verhandlungen guardrail-konform und kommunikativ plausibel führen.",
      0.55, 5.15, 12.2, 0.45, size=12, color=RGBColor(0xA8,0xC4,0xE8), italic=True)

# Three next-steps
rect(sl, 0.35, 5.72, 12.6, 1.42, fill=RGBColor(0x1A,0x23,0x3C))
txbox(sl, "Offene Punkte & Ausblick",
      0.55, 5.77, 4.0, 0.35, size=12, bold=True,
      color=RGBColor(0xD4,0x7E,0x0F))
nexts = [
    "Aktive Walk-Away-Kommunikation implementieren",
    "Dimension-Shift-Kompetenz für nicht-preisliche Konflikte (S11)",
    "Preis-Text-Konsistenz sicherstellen (30,9 % Runden betroffen)",
]
for i, nx in enumerate(nexts):
    txbox(sl, f"{i+1}.  {nx}",
          0.55, 6.15 + i * 0.28, 12.0, 0.28, size=11, color=WHITE)

footer(sl)


# ── Save ──────────────────────────────────────────────────────
out = "evaluation/presentation_tradebridge.pptx"
prs.save(out)
print(f"Saved → {out}")
