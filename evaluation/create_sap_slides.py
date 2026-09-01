"""
TradeBridge 2.0 — SAP Colleague Presentation
~30 min incl. live demo | Idea → Architecture → KPIs → Results → Demo
Run: .venv/bin/python evaluation/create_sap_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy

# ── SAP Brand Palette ──────────────────────────────────────────
SAP_BLUE     = RGBColor(0x00, 0x70, 0xD2)
SAP_DARK     = RGBColor(0x00, 0x3D, 0x7A)
SAP_LIGHT    = RGBColor(0xD6, 0xEA, 0xF8)
SAP_GOLD     = RGBColor(0xF0, 0xA3, 0x00)
GREEN        = RGBColor(0x1A, 0x7A, 0x46)
GREEN_LIGHT  = RGBColor(0xD5, 0xF0, 0xE0)
RED          = RGBColor(0xC0, 0x39, 0x2B)
AMBER        = RGBColor(0xD4, 0x7E, 0x0F)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1A, 0x1A, 0x2E)
GREY         = RGBColor(0x55, 0x55, 0x66)
LIGHT_BG     = RGBColor(0xF5, 0xF7, 0xFA)
DIVIDER      = RGBColor(0xD0, 0xD8, 0xE4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Low-level helpers ──────────────────────────────────────────

def sl():
    return prs.slides.add_slide(BLANK)

def box(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def add_line(tf, text, size=13, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(3)):
    p = tf.add_paragraph()
    p.alignment = align; p.space_before = space_before
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return p

def header(slide, title, subtitle=None, bg=SAP_DARK, h=1.25):
    box(slide, 0, 0, 13.33, h, fill=bg)
    txt(slide, title, 0.45, 0.12, 11, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.82, 11, 0.38,
            size=13, color=RGBColor(0xA8, 0xC8, 0xF0))

def footer(slide, note=None):
    box(slide, 0, 7.18, 13.33, 0.32, fill=SAP_DARK)
    label = note or "TradeBridge 2.0  ·  Tarnbir Singh  ·  DHBW Mannheim / SAP SE  ·  2026"
    txt(slide, label, 0.35, 7.19, 12.5, 0.28, size=9,
        color=RGBColor(0x90, 0xB8, 0xD8))

def slide_number(slide, n):
    txt(slide, str(n), 12.9, 7.19, 0.35, 0.28, size=9,
        color=RGBColor(0x70, 0x98, 0xB8), align=PP_ALIGN.RIGHT)

def pill(slide, text, l, t, color=SAP_BLUE, text_color=WHITE, w=None):
    estimated_w = len(text) * 0.085 + 0.3
    bw = w or estimated_w
    box(slide, l, t, bw, 0.28, fill=color)
    txt(slide, text, l + 0.08, t + 0.03, bw - 0.12, 0.22,
        size=9, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=SAP_DARK)
box(s, 0, 4.2, 13.33, 0.04, fill=SAP_BLUE)

# Brand mark
box(s, 0.55, 0.55, 0.06, 3.4, fill=SAP_BLUE)

txt(s, "TradeBridge 2.0", 0.85, 1.0, 11.5, 1.3,
    size=58, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(s, "Autonome B2B-Verhandlungen mit KI-Agenten",
    0.85, 2.45, 10, 0.6, size=22,
    color=RGBColor(0xA8, 0xC8, 0xF0), align=PP_ALIGN.LEFT)

txt(s, "Bachelorarbeit  ·  DHBW Mannheim / SAP SE",
    0.85, 4.45, 9, 0.4, size=14,
    color=RGBColor(0x80, 0xA8, 0xCC), align=PP_ALIGN.LEFT)
txt(s, "Tarnbir Singh  ·  2026",
    0.85, 4.92, 6, 0.35, size=13,
    color=RGBColor(0x60, 0x88, 0xAA), align=PP_ALIGN.LEFT)

# Agenda pills
pills = ["01  Idee & Problem", "02  Architektur", "03  KPIs", "04  Ergebnisse", "05  Live Demo"]
for i, p_text in enumerate(pills):
    pill(s, p_text, 0.85 + i * 2.45, 5.85,
         color=RGBColor(0x00, 0x50, 0x9A), text_color=RGBColor(0xA8, 0xC8, 0xF0), w=2.25)

footer(s, "SAP SE — Interner Vortrag")
slide_number(s, 1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem / Motivation
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "01  Die Ausgangslage", "Tail-Spend-Verhandlungen im Core Retail — manuell, langsam, teuer")

# Three pain-point cards
pain = [
    ("⏱", "Manueller Aufwand",
     "Einkäufer verhandeln Standardartikel mit geringem Wert manuell — hoher Zeitaufwand, niedriger strategischer ROI."),
    ("📉", "Inkonsistente Ergebnisse",
     "Ohne formalisierte Strategie variieren Einigungspreise stark — je nach Erfahrung und Tagesform des Einkäufers."),
    ("🔒", "Keine Skalierbarkeit",
     "Steigendes Sortimentsvolumen lässt sich nicht mit wachsendem Personal lösen — Automatisierung ist notwendig."),
]
for i, (icon, title, desc) in enumerate(pain):
    l = 0.35 + i * 4.3
    box(s, l, 1.45, 4.05, 3.2, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
    txt(s, icon, l + 0.2, 1.6, 0.7, 0.6, size=28)
    txt(s, title, l + 0.2, 2.25, 3.6, 0.45, size=15, bold=True, color=SAP_DARK)
    tb = s.shapes.add_textbox(Inches(l + 0.2), Inches(2.75), Inches(3.65), Inches(1.7))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = desc
    tf.paragraphs[0].font.size = Pt(11.5)
    tf.paragraphs[0].font.color.rgb = GREY

# Research question
box(s, 0.35, 4.85, 12.63, 1.35, fill=SAP_LIGHT, line_color=SAP_BLUE, line_w=Pt(2))
txt(s, "Forschungsfrage", 0.6, 4.95, 4, 0.32, size=11, bold=True, color=SAP_BLUE)
txt(s,
    "\"Inwiefern laesst sich der Verhandlungsprozess zwischen Lieferant und Haendler im Core Retail"
    " durch ein LLM-basiertes Multi-Agenten-System innerhalb vordefinierter Preis- und Mengenparameter automatisieren?\"",
    0.6, 5.32, 12.2, 0.75, size=13, italic=True, color=SAP_DARK)

footer(s); slide_number(s, 2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — The Idea
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "01  Die Idee", "Zwei KI-Agenten verhandeln autonom — innerhalb menschlich definierter Grenzen")

# Central flow diagram
nodes = [
    (0.35,  2.0, 2.8, 3.0, SAP_DARK,  "🏭",  "Supplier\nAgent",    "Setzt Mindestpreis\n& Strategie"),
    (5.27,  1.7, 2.8, 3.6, SAP_BLUE,  "⚙️",  "Orchestrator\n& ZOPA", "Verhandlungslogik\nConstraint-Check\nHITL-Eskalation"),
    (10.18, 2.0, 2.8, 3.0, RGBColor(0x1E,0x6B,0x4A), "🛒", "Retailer\nAgent", "Setzt Max-Preis\n& Taktik"),
]
for (l, t, w, h, col, icon, title, sub) in nodes:
    box(s, l, t, w, h, fill=col)
    txt(s, icon,  l+w/2-0.25, t+0.2,  0.5, 0.55, size=24, align=PP_ALIGN.CENTER)
    txt(s, title, l+0.1, t+0.85, w-0.15, 0.65, size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub,   l+0.12, t+1.6,  w-0.2,  1.1,  size=10,
        color=RGBColor(0xB8,0xD4,0xF0), align=PP_ALIGN.CENTER)

# Arrows
for ax in [3.25, 8.05]:
    txt(s, "⟵  Angebote  ⟶", ax, 3.1, 1.9, 0.45, size=10,
        color=GREY, align=PP_ALIGN.CENTER)

# Key principles below
principles = [
    ("🔐 Guardrails", "Preisuntergrenze und Mengenkonstraints sind unveränderlich — der Agent kann sie nie unterschreiten."),
    ("🧠 LLM-Strategie", "Claude Sonnet entscheidet Taktik, Konzessionsgröße und Begründung — Runde für Runde."),
    ("👤 HITL-Eskalation", "Menschlicher Eingriff nur bei Stagnation oder Grenzfällen — alles andere läuft autonom."),
    ("📊 Transparenz", "Jede Runde wird mit Reasoning, raw_offer und Taktikklassifikation persistiert."),
]
for i, (title, desc) in enumerate(principles):
    l = 0.35 + (i % 2) * 6.5
    t = 5.25 + (i // 2) * 0.85
    box(s, l, t, 6.2, 0.72, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
    txt(s, title, l+0.15, t+0.08, 2.0, 0.3, size=11, bold=True, color=SAP_DARK)
    txt(s, desc,  l+0.15, t+0.38, 5.9, 0.3, size=10, color=GREY)

footer(s); slide_number(s, 3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Deep-Dive
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "02  Systemarchitektur", "Modularer Aufbau — jede Komponente hat eine klar abgegrenzte Verantwortung")

# Left: component stack
components = [
    (SAP_DARK,  "FastAPI Backend",         "REST-API, Session-Management, ZOPA-Berechnung"),
    (SAP_BLUE,  "LLM-Pipeline (4 Stufen)", "Tactic → Aspiration → Risk → Offer — sequenziell, prompt-konfiguriert"),
    (RGBColor(0x1E,0x6B,0x4A), "Constraint Validator + Retry", "Validierung auf raw_offer; max. 3 Re-Prompts vor Clamping"),
    (RGBColor(0x6A,0x1B,0x9A), "Opponent Model",              "Boulware/Conceder-Klassifikation nach Faratin et al. (1998)"),
    (RGBColor(0x8C,0x4A,0x00), "Evaluationsmodul",            "CSR / WAA / ZU / BP — automatisiert über 140 Sitzungen"),
]
for i, (col, title, desc) in enumerate(components):
    t = 1.45 + i * 1.06
    box(s, 0.35, t, 0.08, 0.75, fill=col)
    box(s, 0.55, t, 5.8, 0.82, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
    txt(s, title, 0.72, t+0.06, 5.4, 0.35, size=12, bold=True, color=col)
    txt(s, desc,  0.72, t+0.42, 5.4, 0.35, size=10, color=GREY)

# Right: tech stack
box(s, 6.7, 1.45, 6.3, 5.65, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
txt(s, "Tech Stack", 6.9, 1.55, 5.8, 0.35, size=12, bold=True, color=SAP_BLUE)

stack = [
    ("LLM",         "Claude Sonnet 4.x  ·  Temperature = 0"),
    ("Backend",     "Python 3.13  ·  FastAPI  ·  Pydantic v2"),
    ("Frontend",    "React 19  ·  TypeScript  ·  Tailwind CSS"),
    ("Protokoll",   "A2A (Agent-to-Agent)  ·  REST-basiert"),
    ("Evaluation",  "140 Sitzungen  ·  14 Szenarien  ·  10 Läufe"),
    ("Persistenz",  "JSON-Session-Store  ·  raw_offer-Feld"),
]
for i, (k, v) in enumerate(stack):
    t = 2.05 + i * 0.82
    box(s, 6.8, t, 1.5, 0.62, fill=SAP_LIGHT)
    txt(s, k, 6.88, t+0.14, 1.35, 0.3, size=10, bold=True, color=SAP_DARK, align=PP_ALIGN.CENTER)
    txt(s, v, 8.45, t+0.14, 4.35, 0.55, size=11, color=DARK)

footer(s); slide_number(s, 4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — KPIs
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "03  Evaluation — vier KPIs", "Hybrides Design: quantitativ automatisiert + qualitativ Human-as-a-Judge")

kpi_data = [
    (SAP_DARK,  "CSR",
     "Constraint Satisfaction Rate",
     "Hält der Agent Preis- und Mengenlimits aus eigener Kraft ein?",
     "Gemessen auf raw_offer — erstem LLM-Output vor Korrektur.",
     "≥ 95 %"),
    (SAP_BLUE,  "WAA",
     "Walk-Away Accuracy",
     "Trifft der Agent die richtige Deal-/Abbruch-Entscheidung?",
     "F₁ auf ZOPA/No-ZOPA-Klassifikation. FAR = 0 ist hartes Kriterium.",
     "F₁ ≥ 0,90\nFAR = 0,0"),
    (GREEN,     "ZU",
     "ZOPA Utilization",
     "Wo landet die Einigung im Verhandlungskorridor?",
     "0 = Supplier-Maximum, 1 = Supplier-Minimum. Diagnostisch, kein SW.",
     "Diagnostisch"),
    (RGBColor(0x6A,0x1B,0x9A), "BP",
     "Business Plausibility",
     "Klingt das nach echter B2B-Verhandlung?",
     "Retail-Experte bewertet Argumentationsqualität, Konzessionslogik, Reaktionsadäquatheit.",
     "Human-as-a-Judge"),
]

for i, (col, short, name, q, method, crit) in enumerate(kpi_data):
    l = 0.35 + (i % 2) * 6.5
    t = 1.45 + (i // 2) * 2.6
    box(s, l, t, 6.2, 2.4, fill=WHITE, line_color=col, line_w=Pt(2.5))
    box(s, l, t, 0.9, 2.4, fill=col)
    txt(s, short, l+0.06, t+0.78, 0.78, 0.7, size=22, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name,   l+1.05, t+0.08, 5.05, 0.4,  size=13, bold=True, color=col)
    txt(s, q,      l+1.05, t+0.52, 5.05, 0.45, size=11, color=DARK)
    txt(s, method, l+1.05, t+1.0,  5.05, 0.5,  size=10, color=GREY, italic=True)
    box(s, l+1.05, t+1.62, 4.8, 0.55, fill=SAP_LIGHT)
    txt(s, "Kriterium:  " + crit, l+1.2, t+1.72, 4.5, 0.38,
        size=10, bold=True, color=SAP_DARK)

footer(s); slide_number(s, 5)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Results Overview (all 4 KPIs)
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "04  Ergebnisse — Übersicht",
       "140 Sitzungen · 14 Szenarien · 10 Läufe · Claude Sonnet · Temperature = 0")

# Big 4 metric cards
metrics = [
    (SAP_DARK,  "CSR",       "98,4 %",  "± 1,6 %", "Kriterium ≥ 95 %", "✓ erfüllt", GREEN),
    (SAP_BLUE,  "WAA  F₁",   "0,933",   "",         "Kriterium ≥ 0,90", "✓ erfüllt", GREEN),
    (SAP_BLUE,  "WAA  FAR",  "0,000",   "",         "Hartes Kriterium", "✓ erfüllt", GREEN),
    (RGBColor(0x6A,0x1B,0x9A), "BP", "8 / 14", "überwiegend / hoch", "Human-as-a-Judge", "✓ positiv", GREEN),
]
for i, (col, label, val, sub, crit, verdict, vcol) in enumerate(metrics):
    l = 0.35 + i * 3.25
    box(s, l, 1.45, 3.1, 2.55, fill=SAP_DARK)
    box(s, l, 1.45, 3.1, 0.06, fill=col)
    txt(s, label,   l+0.15, 1.55, 2.8, 0.38, size=12, color=RGBColor(0x90,0xB8,0xD8))
    txt(s, val,     l+0.12, 1.93, 2.86, 0.78, size=30, bold=True,
        color=vcol, align=PP_ALIGN.CENTER)
    if sub:
        txt(s, sub, l+0.12, 2.68, 2.86, 0.3, size=10,
            color=RGBColor(0x80,0xA0,0xC0), align=PP_ALIGN.CENTER)
    txt(s, crit,    l+0.12, 3.0, 2.86, 0.32, size=9, italic=True,
        color=RGBColor(0x70,0x90,0xB0), align=PP_ALIGN.CENTER)
    txt(s, verdict, l+0.12, 3.3, 2.86, 0.38, size=11, bold=True,
        color=vcol, align=PP_ALIGN.CENTER)

# Key insight boxes
insights = [
    (GREEN,  "✓",  "Null BATNA-Verletzungen",
     "In keiner einzigen der 140 Sitzungen schließt das System einen Deal, der den\n"
     "Reservationspreis des Agenten unterschreitet. FAR = 0,0 über alle Läufe."),
    (AMBER,  "⚠",  "Narrow-ZOPA bleibt herausfordernd",
     "10 verlorene Deals (FN) trotz vorhandener ZOPA — alle in engen Preisfenstern\n"
     "(S04–S06). 1× autonomer Walk-Away obwohl ZOPA vorhanden (S06, Lauf 8)."),
    (SAP_BLUE, "✓", "Volume-Leverage funktioniert",
     "In allen 10 Läufen setzt der Retailer-Agent den Mengenhebel korrekt ein (S13).\n"
     "ZU = 0,618 — einziges Szenario mit Händler-Dominanz."),
]
for i, (col, icon, title, desc) in enumerate(insights):
    l = 0.35 + i * 4.32
    box(s, l, 4.2, 4.1, 2.65, fill=WHITE, line_color=col, line_w=Pt(2))
    txt(s, icon,  l+0.15, 4.3,  0.38, 0.45, size=18, bold=True, color=col)
    txt(s, title, l+0.6,  4.35, 3.35, 0.4,  size=12, bold=True, color=DARK)
    tb = s.shapes.add_textbox(Inches(l+0.15), Inches(4.82),
                               Inches(3.8), Inches(1.85))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = desc
    tf.paragraphs[0].font.size = Pt(10.5)
    tf.paragraphs[0].font.color.rgb = GREY

footer(s); slide_number(s, 6)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — WAA Detail (Confusion Matrix)
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "04  Walk-Away Accuracy — Detail",
       "Positive Klasse: Deal bei vorhandener ZOPA  ·  N = 120 klassifizierte Sitzungen")

# Confusion matrix
box(s, 0.35, 1.45, 6.2, 3.5, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
txt(s, "Aggregierte Konfusionsmatrix (10 Läufe, S11/S12 ausgeschlossen)",
    0.5, 1.52, 5.9, 0.35, size=11, bold=True, color=SAP_BLUE)

# Headers
txt(s, "ZOPA vorhanden", 2.0, 2.0, 2.1, 0.3, size=10, bold=True,
    color=SAP_DARK, align=PP_ALIGN.CENTER)
txt(s, "Keine ZOPA", 4.15, 2.0, 2.0, 0.3, size=10, bold=True,
    color=SAP_DARK, align=PP_ALIGN.CENTER)
txt(s, "Deal", 0.42, 2.5, 1.5, 0.55, size=13, bold=True,
    color=SAP_DARK, align=PP_ALIGN.RIGHT)
txt(s, "Abbruch", 0.42, 3.3, 1.5, 0.55, size=13, bold=True,
    color=SAP_DARK, align=PP_ALIGN.RIGHT)

cells = [
    (2.0,  2.4, "TP = 70",  GREEN,       "Korrekte Deals"),
    (4.15, 2.4, "FP = 0",   GREEN,       "0 BATNA-Fehler"),
    (2.0,  3.2, "FN = 10",  AMBER,       "Verlorene Deals"),
    (4.15, 3.2, "TN = 40",  GREEN,       "Korrekte Abbrüche"),
]
for (cl, ct, val, col, sub) in cells:
    box(s, cl, ct, 1.95, 0.68, fill=col)
    txt(s, val, cl+0.08, ct+0.08, 1.8, 0.38, size=18, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub, cl+0.08, ct+0.44, 1.8, 0.22, size=8,
        color=WHITE, align=PP_ALIGN.CENTER)

# Derived metrics
box(s, 0.35, 5.1, 6.2, 1.65, fill=SAP_LIGHT)
txt(s, "Precision = 1,000   ·   Recall = 0,875   ·   F₁ = 0,933",
    0.55, 5.22, 5.8, 0.38, size=13, bold=True, color=SAP_DARK)
txt(s, "FAR = 0,000 (= FP / (FP+TN))   ·   FWR = 0,125 (= FN / (FN+TP))",
    0.55, 5.62, 5.8, 0.35, size=12, color=SAP_DARK)
txt(s, "FAR = 0 bedeutet: keine einzige BATNA-Verletzung in 140 Sitzungen",
    0.55, 6.0, 5.8, 0.35, size=11, italic=True, color=GREY)

# Right: scenario breakdown of FN
box(s, 6.8, 1.45, 6.15, 5.3, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
txt(s, "FN-Analyse — 10 verlorene Deals", 6.95, 1.52, 5.8, 0.35,
    size=11, bold=True, color=SAP_BLUE)

fn_data = [
    ("S04  Makita Schlagbohrer",     "Narrow ZOPA €6",  2, AMBER),
    ("S05  Bosch X-LOCK Scheiben",   "Narrow ZOPA €1,5",3, RED),
    ("S06  GARDENA Water Control",   "Narrow ZOPA €5",  4, RED),
    ("S06  Lauf 8 — Sonderfall",     "Autonomer Walk-Away (failed)",1, RED),
    ("S13  Kärcher Schaum-Set",      "Volume Leverage", 1, AMBER),
]
for i, (name, reason, n, col) in enumerate(fn_data):
    t = 2.05 + i * 0.88
    box(s, 6.9, t, 0.45, 0.62, fill=col)
    txt(s, str(n), 6.9, t+0.1, 0.45, 0.4, size=16, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, name,   7.45, t+0.05, 5.35, 0.3,  size=11, bold=True, color=DARK)
    txt(s, reason, 7.45, t+0.35, 5.35, 0.25, size=10, color=GREY, italic=True)

txt(s, "Alle 10 FN-Fälle: Verhandlung korrekt abgebrochen, aber Einigung wäre möglich gewesen.",
    6.9, 6.62, 5.9, 0.38, size=9, color=GREY, italic=True)

footer(s); slide_number(s, 7)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — ZU + BP Summary
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "04  ZU & Business Plausibility",
       "Einigungsposition im Korridor + qualitative Verhandlungsbeurteilung")

# ZU section
box(s, 0.35, 1.45, 6.3, 5.65, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
txt(s, "ZOPA Utilization — Ø 0,303 ± 0,043",
    0.5, 1.52, 6.0, 0.38, size=12, bold=True, color=SAP_BLUE)
txt(s, "67,5 % lieferantendominant  ·  16,9 % ausgeglichen  ·  15,7 % händlerdominant",
    0.5, 1.94, 5.9, 0.3, size=10, color=GREY, italic=True)

zu_rows = [
    ("S01–S03", "Wide ZOPA",       [0.219, 0.160, 0.312], AMBER),
    ("S04–S06", "Narrow ZOPA",     [0.094, 0.286, 0.183], AMBER),
    ("S11–S12", "Asymmetrisch",    [0.293, 0.279],         AMBER),
    ("S13",     "Volume Leverage", [0.618],                SAP_BLUE),
    ("S14",     "Volume Leverage", [0.517],                GREEN),
]
bar_start = 2.2; bar_max = 3.6
for i, (sid, cat, vals, col) in enumerate(zu_rows):
    t = 2.42 + i * 0.82
    mean = sum(vals) / len(vals)
    txt(s, sid,   0.5,  t+0.05, 0.85, 0.35, size=9,  bold=True, color=DARK)
    txt(s, cat,   1.38, t+0.05, 1.75, 0.35, size=9,  color=GREY)
    bar_w = bar_max * mean
    box(s, bar_start, t+0.05, bar_w, 0.32, fill=col)
    txt(s, f"{mean:.3f}", bar_start + bar_w + 0.1, t+0.07, 0.65, 0.28, size=9, color=DARK)
    if col == SAP_BLUE:
        txt(s, "← Hebel wirkt", bar_start + bar_w + 0.8, t+0.07, 1.2, 0.28,
            size=9, bold=True, color=SAP_BLUE)

# ZU interpretation
box(s, 0.35, 6.25, 6.3, 0.85, fill=SAP_LIGHT)
txt(s, "Systematischer Supplier-Bias durch Ankerstrategie im Prompt-Design — korrigierbar durch\n"
    "Anpassung des `initial_anchor`-Parameters in der Agentenpersönlichkeit.",
    0.5, 6.33, 6.0, 0.68, size=10, color=SAP_DARK, italic=True)

# BP section
box(s, 6.9, 1.45, 6.1, 5.65, fill=WHITE, line_color=DIVIDER, line_w=Pt(1))
txt(s, "Business Plausibility — Einstufungen",
    7.05, 1.52, 5.8, 0.35, size=12, bold=True, color=SAP_BLUE)

bp_items = [
    (GREEN,      "Hoch plausibel",           1,  "S13 — Volume Leverage"),
    (SAP_BLUE,   "Überwiegend plausibel",     7,  "S01–S04, S12, S14"),
    (AMBER,      "Eingeschränkt plausibel",   6,  "S05–S11"),
    (DIVIDER,    "Nicht plausibel",           0,  "—"),
]
for i, (col, label, n, examples) in enumerate(bp_items):
    t = 2.05 + i * 1.2
    box(s, 7.0, t, 0.8, 0.9, fill=col)
    txt(s, str(n), 7.0, t+0.18, 0.8, 0.5, size=24, bold=True,
        color=WHITE if col != DIVIDER else GREY, align=PP_ALIGN.CENTER)
    txt(s, label,    7.95, t+0.06, 4.9, 0.36, size=12, bold=True, color=DARK)
    txt(s, examples, 7.95, t+0.48, 4.9, 0.3,  size=10, color=GREY, italic=True)

box(s, 6.9, 6.25, 6.1, 0.85, fill=SAP_LIGHT)
txt(s, "Kein Szenario als 'nicht plausibel' eingestuft. Hauptschwäche:\n"
    "fehlende Walk-Away-Kommunikation in No-ZOPA-Szenarien (S07–S10).",
    7.05, 6.33, 5.8, 0.68, size=10, color=SAP_DARK, italic=True)

footer(s); slide_number(s, 8)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Live Demo Cue
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=SAP_DARK)
box(s, 0, 0, 0.12, 7.5, fill=SAP_BLUE)

txt(s, "05", 1.0, 0.6, 2.0, 1.0, size=60, bold=True,
    color=RGBColor(0x00, 0x50, 0x9A))
txt(s, "Live Demo", 1.0, 1.65, 11, 1.0, size=52, bold=True, color=WHITE)
txt(s, "TradeBridge 2.0 — End-to-End Verhandlung",
    1.0, 2.75, 10, 0.5, size=18,
    color=RGBColor(0xA8, 0xC8, 0xF0))

box(s, 1.0, 3.5, 11.0, 2.5, fill=RGBColor(0x00, 0x28, 0x5A))
txt(s, "Demo-Ablauf", 1.2, 3.62, 4, 0.38, size=13, bold=True,
    color=SAP_GOLD)
demo_steps = [
    "1.  Supplier-Dashboard → Produkt auswählen → Proaktives Angebot mit privaten Limits",
    "2.  Retailer-Dashboard → Angebot im Posteingang → Eigene Max-Grenze setzen → Auto-Verhandlung starten",
    "3.  Live: KI-Agenten verhandeln Runde für Runde mit Taktik & Begründung",
    "4.  HITL-Eskalation oder Deal-Bestätigung durch beide Parteien",
]
tb = s.shapes.add_textbox(Inches(1.2), Inches(4.1), Inches(10.6), Inches(1.8))
tb.word_wrap = True
tf = tb.text_frame; tf.word_wrap = True
for i, step in enumerate(demo_steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = step; p.space_before = Pt(3)
    p.font.size = Pt(12); p.font.color.rgb = WHITE

txt(s, "localhost:5173  ·  Backend: localhost:8002",
    1.0, 6.2, 8, 0.38, size=12,
    color=RGBColor(0x70, 0x98, 0xB8), italic=True)

footer(s, "TradeBridge 2.0  ·  Live Demo")
slide_number(s, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Summary & Outlook
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=LIGHT_BG)
header(s, "Zusammenfassung & Ausblick",
       "Was wurde gezeigt — und was kommt als nächstes?")

# Left: what was achieved
box(s, 0.35, 1.45, 5.95, 5.55, fill=WHITE, line_color=GREEN, line_w=Pt(2))
txt(s, "✓  Erreicht", 0.55, 1.55, 5.5, 0.38, size=13, bold=True, color=GREEN)

achieved = [
    "Alle 3 quantitativen Erfolgskriterien erfüllt",
    "CSR 98,4 % — Guardrails hochzuverlässig",
    "FAR = 0,0 — null BATNA-Verletzungen in 140 Sitzungen",
    "F₁ = 0,933 — ausgewogene Walk-Away-Entscheidung",
    "Volume-Leverage in 10/10 Läufen korrekt eingesetzt",
    "Kommunikativ plausibles Verhalten in 8/14 Szenarien",
    "End-to-End-System: Frontend → API → LLM-Agenten",
]
tb = s.shapes.add_textbox(Inches(0.55), Inches(2.0), Inches(5.6), Inches(4.6))
tb.word_wrap = True
tf = tb.text_frame; tf.word_wrap = True
for i, a in enumerate(achieved):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run(); r.text = f"  ·  {a}"
    r.font.size = Pt(11.5); r.font.color.rgb = DARK

# Right: limitations & next steps
box(s, 6.55, 1.45, 6.43, 2.55, fill=WHITE, line_color=AMBER, line_w=Pt(2))
txt(s, "⚠  Bekannte Grenzen", 6.75, 1.55, 6.0, 0.38, size=13, bold=True, color=AMBER)
limits = [
    "Narrow ZOPA: 20–40 % Fehlquote (Anker-Varianz)",
    "Keine aktive Walk-Away-Kommunikation implementiert",
    "Preis-Text-Diskrepanz: 30,9 % der Runden betroffen",
    "S11: kein Dimension-Shift bei Nicht-Preis-Konflikten",
]
tb2 = s.shapes.add_textbox(Inches(6.75), Inches(2.0), Inches(6.1), Inches(1.75))
tb2.word_wrap = True
tf2 = tb2.text_frame; tf2.word_wrap = True
for i, lim in enumerate(limits):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run(); r.text = f"  ·  {lim}"
    r.font.size = Pt(11); r.font.color.rgb = DARK

box(s, 6.55, 4.2, 6.43, 2.8, fill=SAP_LIGHT, line_color=SAP_BLUE, line_w=Pt(2))
txt(s, "→  Nächste Schritte", 6.75, 4.3, 5.8, 0.38, size=13, bold=True, color=SAP_BLUE)
nexts = [
    "Aktive Abbruchkommunikation (Walk-Away-Signal)",
    "Regelbasierter Ankerpreis für enge ZOPAs",
    "Domänenspezifische Prompt-Erweiterung (Saisonalität, Margen)",
    "Preis-Text-Konsistenzprüfung im Pipeline-Design",
    "Produktive Integration in SAP Ariba / S/4HANA",
]
tb3 = s.shapes.add_textbox(Inches(6.75), Inches(4.75), Inches(6.1), Inches(2.1))
tb3.word_wrap = True
tf3 = tb3.text_frame; tf3.word_wrap = True
for i, nx in enumerate(nexts):
    p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run(); r.text = f"  {i+1}.  {nx}"
    r.font.size = Pt(11); r.font.color.rgb = DARK

footer(s); slide_number(s, 10)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Closing / Thank You
# ═══════════════════════════════════════════════════════════════
s = sl()
box(s, 0, 0, 13.33, 7.5, fill=SAP_DARK)
box(s, 0, 0, 0.12, 7.5, fill=SAP_BLUE)
box(s, 0, 3.72, 13.33, 0.04, fill=SAP_BLUE)

txt(s, "Danke!", 1.0, 1.0, 11, 1.5,
    size=64, bold=True, color=WHITE)
txt(s, "Fragen & Diskussion",
    1.0, 2.65, 10, 0.55, size=22,
    color=RGBColor(0xA8, 0xC8, 0xF0))

txt(s, "Tarnbir Singh", 1.0, 3.95, 8, 0.45, size=16, bold=True, color=WHITE)
txt(s, "DHBW Mannheim  ·  SAP SE  ·  2026",
    1.0, 4.45, 8, 0.38, size=13,
    color=RGBColor(0x80, 0xA8, 0xCC))

txt(s, "GitHub / Code:  Agentic-2.0-Retail-Industry",
    1.0, 5.3, 9, 0.38, size=12,
    color=RGBColor(0x60, 0x88, 0xAA), italic=True)
txt(s, "Backend:  localhost:8002  ·  Frontend:  localhost:5173",
    1.0, 5.72, 9, 0.35, size=11,
    color=RGBColor(0x50, 0x78, 0x9A), italic=True)

footer(s, "TradeBridge 2.0  ·  Tarnbir Singh  ·  DHBW Mannheim / SAP SE  ·  2026")
slide_number(s, 11)


# ── Save ──────────────────────────────────────────────────────
out = "evaluation/presentation_sap.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
